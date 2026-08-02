from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

BASE = Path('/home/pgupta11/Projects/STAT6_PPI/PPI_program_writing/R_group_mapping/advanced_problems')
REPORT = BASE / 'VS_PPI_Leo_Bicyclic_headgroup_screen_06052026_report.html'
OUT = BASE / 'VS_PPI_Leo_Bicyclic_headgroup_screen_06052026_presentation_15slides.pptx'

IMG_INTRO = BASE / 'report_assets' / 'intro_slide_rgroup_report_1920x1080.png'
IMG_SCAFF = BASE / 'scaffold_example.png'
IMG_PROP = BASE / 'plotly_2d_corr.png'
IMG_GREEN = BASE / 'green_highlight.png'
IMG_RED = BASE / 'red_highlight.png'

TITLE_COLOR = RGBColor(17, 34, 51)
ACCENT = RGBColor(11, 110, 79)


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    tf = slide.shapes.title.text_frame
    tf.paragraphs[0].font.color.rgb = TITLE_COLOR
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(40)
    slide.placeholders[1].text = subtitle
    stf = slide.placeholders[1].text_frame
    stf.paragraphs[0].font.size = Pt(20)
    stf.paragraphs[0].font.color.rgb = RGBColor(60, 80, 100)
    return slide


def add_bullets_slide(prs, title, bullets, level2=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True

    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(40, 55, 72)
        if level2 and i in level2:
            for sub in level2[i]:
                sp = body.add_paragraph()
                sp.text = sub
                sp.level = 1
                sp.font.size = Pt(18)
                sp.font.color.rgb = RGBColor(70, 90, 110)
    return slide


def add_image_slide(prs, title, image_path, caption):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True

    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(0.6), Inches(1.0), width=Inches(12.1), height=Inches(5.9))
    tx = slide.shapes.add_textbox(Inches(0.7), Inches(6.2), Inches(12.0), Inches(0.6))
    p = tx.text_frame.paragraphs[0]
    p.text = caption
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(60, 80, 100)
    return slide


def add_two_image_compare_slide(prs, title, left_img, right_img, left_label, right_label):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True

    if left_img.exists():
        slide.shapes.add_picture(str(left_img), Inches(0.7), Inches(1.2), width=Inches(5.8), height=Inches(4.2))
    if right_img.exists():
        slide.shapes.add_picture(str(right_img), Inches(6.8), Inches(1.2), width=Inches(5.8), height=Inches(4.2))

    ltx = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(5.6), Inches(0.7))
    ltp = ltx.text_frame.paragraphs[0]
    ltp.text = left_label
    ltp.font.size = Pt(16)
    ltp.font.color.rgb = ACCENT
    ltp.font.bold = True

    rtx = slide.shapes.add_textbox(Inches(6.9), Inches(5.5), Inches(5.6), Inches(0.7))
    rtp = rtx.text_frame.paragraphs[0]
    rtp.text = right_label
    rtp.font.size = Pt(16)
    rtp.font.color.rgb = RGBColor(190, 70, 30)
    rtp.font.bold = True

    return slide


def build():
    prs = Presentation()

    # 1
    add_title_slide(
        prs,
        'VS_PPI Leo Bicyclic Headgroup Screen',
        '15-slide review of HTML report: scaffold insights, deep dive, docking visualizer, and filtering workflows'
    )

    # 2
    add_image_slide(
        prs,
        'Report At A Glance',
        IMG_INTRO,
        'Interactive single-file HTML report generated from docking poses + interaction-count CSV.'
    )

    # 3
    add_bullets_slide(
        prs,
        'Input and Processing Context',
        [
            'Inputs: docking pose SDF + interaction CSV + optional property annotations',
            'Pipeline clusters molecules by scaffold and substitution pattern',
            'Outputs: Central Ideas, Scaffold Deep Dive, and 3D docking pose review',
            'This report file is large and self-contained (~91 MB HTML)'
        ]
    )

    # 4
    add_bullets_slide(
        prs,
        'HTML Structure Reviewed',
        [
            'Main sections detected in HTML:',
            'R-group Docking Insight Report',
            'Hydrogen Bonding Residues',
            'Molecule Properties',
            'Central Ideas',
            'Scaffold Deep Dive (Central Ideas)'
        ]
    )

    # 5
    add_bullets_slide(
        prs,
        'Scale and Coverage',
        [
            'Scaffolds shown in report header: 732',
            'Central Ideas shown in report header: 146',
            'Large embedded data and interactive JS components',
            'Designed for medicinal chemistry triage and rapid review'
        ]
    )

    # 6
    add_image_slide(
        prs,
        'Scaffold Panel: Why It Matters',
        IMG_SCAFF,
        'Users can rapidly compare substitution patterns and prioritize scaffold families before deep-dive analysis.'
    )

    # 7
    add_bullets_slide(
        prs,
        'Central Ideas Workflow',
        [
            'Central Ideas summarize top scaffold families for first-pass decision making',
            'Cards expose member count, score trends, and interaction-driven prioritization',
            'Panel click behavior opens deep-dive details and docking pose visualizer context',
            'Supports iterative hit triage by scaffold class'
        ]
    )

    # 8
    add_bullets_slide(
        prs,
        'Scaffold Deep Dive Behavior',
        [
            'Deep Dive includes Central Idea scaffolds with member_count >= 3',
            'Users inspect molecule-level examples within each scaffold family',
            'Per-scaffold drill-down supports substitution strategy selection',
            'Designed to bridge scaffold-level and pose-level evidence'
        ]
    )

    # 9
    add_bullets_slide(
        prs,
        'Docking Pose Visualizer Window',
        [
            'Each panel can launch an interactive 3D docking pose visualizer',
            'Viewer shows protein environment and docked ligand geometry',
            'Supports rotation, zoom, pan, and multi-ligand overlay comparisons',
            'Enables fast visual validation of scaffold hypotheses'
        ]
    )

    # 10
    add_bullets_slide(
        prs,
        'Substructure Search and Motif Controls',
        [
            'Substructure search supports SMARTS/SMILES query modes',
            'Exact match and motif-exclusion workflows are available',
            'Filtering can remove unwanted chemotypes before prioritization',
            'Useful for focusing on tractable medicinal chemistry space'
        ]
    )

    # 11
    add_two_image_compare_slide(
        prs,
        'Hydrogen Bond Residue Filtering',
        IMG_GREEN,
        IMG_RED,
        'Pass: scaffolds matching selected donor/acceptor residue logic',
        'Fail: scaffolds excluded by residue-based interaction criteria'
    )

    # 12
    add_image_slide(
        prs,
        'Molecule Properties and ADME Analytics',
        IMG_PROP,
        'Property panel supports ADME and physico-chemical analysis with histogram/filter and correlation views.'
    )

    # 13
    add_bullets_slide(
        prs,
        'Custom Property Extensibility',
        [
            'Report can list precomputed properties (e.g., molecular weight) and ML-generated ADME fields',
            'Property panel supports custom columns and range filtering',
            'Interaction Count is exposed as a sortable/filterable property',
            'Workflow adapts to new project-specific descriptors without UI redesign'
        ]
    )

    # 14
    add_bullets_slide(
        prs,
        'Download and Export Capabilities',
        [
            'Users can download all poses for a scaffold series',
            'Users can also export selected individual molecules from deep-dive',
            'Supports practical handoff to modeling, synthesis, and project tracking',
            'Enables rapid shortlist creation for follow-up experiments'
        ]
    )

    # 15
    add_bullets_slide(
        prs,
        'Key Takeaways and Recommended Use',
        [
            'Use Central Ideas for fast scaffold triage',
            'Use Deep Dive + 3D visualizer for pose-level confidence checks',
            'Use substructure/motif/H-bond filters to narrow chemistry space',
            'Use property + ADME analytics to prioritize developable candidates',
            'Use export tools to finalize scaffold and molecule shortlists'
        ]
    )

    prs.save(str(OUT))
    print(OUT)


if __name__ == '__main__':
    build()
